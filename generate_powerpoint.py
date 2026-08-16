import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette - Microsoft Fabric Dark Obsidian / Cyan / Slate
    BG_DARK = RGBColor(11, 19, 43)        # #0B132B Dark Obsidian
    CARD_BG = RGBColor(28, 37, 65)        # #1C2541 Slate Navy
    ACCENT_CYAN = RGBColor(72, 202, 228)   # #48CAE4 Bright Cyan
    ACCENT_BLUE = RGBColor(0, 119, 182)    # #0077B6 Deep Blue
    TEXT_WHITE = RGBColor(255, 255, 255)   # White
    TEXT_MUTED = RGBColor(200, 215, 230)  # Soft Muted Gray/Blue
    ACCENT_GREEN = RGBColor(74, 222, 128)  # Emerald Green
    ACCENT_GOLD = RGBColor(251, 191, 36)   # Amber Gold

    def set_background(slide, color=BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="EV FLEET ANALYTICS PLATFORM"):
        # Header Banner / Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        p_cat.font.name = "Arial"

        # Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, border_color=ACCENT_CYAN, card_bg=CARD_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = card_bg
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = border_color
        p_title.font.name = "Arial"
        p_title.space_after = Pt(10)

        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_MUTED
            p.font.name = "Arial"
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    # Accent decorative box
    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    dec.fill.solid()
    dec.fill.fore_color.rgb = ACCENT_CYAN
    dec.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "ELECTRIC VEHICLE (EV) FLEET CHARGING & BATTERY HEALTH ANALYTICS PLATFORM"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Arial"
    p1.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "An Enterprise-Grade Solution for Data Engineering, Cloud Warehousing, RPA Orchestration, Business Intelligence & AI Insights"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_CYAN
    p2.font.name = "Arial"
    p2.space_after = Pt(28)

    p3 = tf1.add_paragraph()
    p3.text = "Technology Stack: Snowflake | Alteryx | UiPath RPA | Power BI & MS Fabric | Node.js | Chart.js & AI Engine"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Business Problem
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Executive Summary & Business ROI")

    add_card(slide2, 0.8, 1.8, 5.6, 5.0, "Operational Challenges", [
        "Unmonitored Battery Degradation leading to sudden fleet vehicle failure.",
        "Uncontrolled Charging Expenses during high-cost peak grid demand periods.",
        "Data Silos across telemetry, location, maintenance logs, and financial records.",
        "Lack of Predictive Insights to schedule timely battery maintenance before breakdown."
    ], border_color=RGBColor(239, 68, 68))

    add_card(slide2, 6.8, 1.8, 5.6, 5.0, "Delivered Platform Value & ROI", [
        "14% Monthly Energy Cost Savings via AI-driven peak-shifting recommendations.",
        "Continuous SOH Tracking to maintain battery State of Health > 80%.",
        "Automated Data Pipelines delivering real-time telemetry to Snowflake Data Warehouse.",
        "6 Interactive Fabric Dashboards + AI Assistant Natural Language Console."
    ], border_color=ACCENT_GREEN)

    # -------------------------------------------------------------
    # SLIDE 3: Architecture & Data Flow
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "End-to-End Enterprise Architecture")

    add_card(slide3, 0.8, 1.8, 3.6, 5.0, "1. Ingestion & Data Prep", [
        "Raw Kaggle EV telemetry CSV data",
        "Alteryx ETL workflow cleans data",
        "Handles missing values & duplicates",
        "Calculates kWh/100km efficiency"
    ], border_color=ACCENT_CYAN)

    add_card(slide3, 4.8, 1.8, 3.6, 5.0, "2. Warehouse & Automation", [
        "Snowflake Star Schema (RAW/CORE/MART)",
        "CDC Streams & 60-min Scheduled Tasks",
        "UiPath RPA orchestrates SP execution",
        "Automates PBI refresh & executive emails"
    ], border_color=ACCENT_BLUE)

    add_card(slide3, 8.8, 1.8, 3.6, 5.0, "3. BI & AI Insights UI", [
        "6 Dark-Blue Microsoft Fabric pages",
        "Production DAX KPI measures",
        "Chart.js interactive web app",
        "AI SOH & cost prediction engine"
    ], border_color=ACCENT_GOLD)

    # -------------------------------------------------------------
    # SLIDE 4: Alteryx ETL Pipeline
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "Data Cleansing & Transformation (Alteryx)")

    add_card(slide4, 0.8, 1.8, 5.6, 5.0, "ETL Pipeline Modules", [
        "Input Data Tool: Loads raw vehicle telemetry & session data.",
        "Data Cleansing & Select: Normalizes strings, handles white space.",
        "Imputation: Resolves missing ambient temp & trip distance.",
        "Unique & Filter: Eliminates duplicates on composite key (Vehicle_ID, Start_Time)."
    ], border_color=ACCENT_CYAN)

    add_card(slide4, 6.8, 1.8, 5.6, 5.0, "Feature Engineering", [
        "Energy Efficiency: Computes (Energy_Consumed / Distance) * 100.",
        "SOH Degradation Bucket: Classifies Battery SOH into Optimal (>90%), Moderate (80-90%), Alert (<80%).",
        "Cost Calculation: Applies peak/off-peak rate tariffs to session energy.",
        "Cleaned CSV Export: Formats clean data for Snowflake ingestion stage."
    ], border_color=ACCENT_CYAN)

    # -------------------------------------------------------------
    # SLIDE 5: Snowflake Cloud Data Warehouse
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Cloud Data Warehousing (Snowflake)")

    add_card(slide5, 0.8, 1.8, 5.6, 5.0, "Multi-Layered Architecture", [
        "RAW Schema: Staging landing tables for raw session loads.",
        "CORE Schema: Production Star Schema DDL.",
        "• DIM_VEHICLE (Vehicle_SK, Model, Capacity, Fleet_ID)",
        "• FACT_CHARGING_SESSIONS (Session_SK, SOH, kWh, Cost, Efficiency)",
        "MART Schema: Pre-computed business analytics views."
    ], border_color=ACCENT_BLUE)

    add_card(slide5, 6.8, 1.8, 5.6, 5.0, "CDC & Stored Procedures", [
        "SP_POPULATE_STAR_SCHEMA(): Automates ETL from RAW into CORE Star Schema.",
        "SP_CALCULATE_BATTERY_DEGRADATION(): Computes monthly degradation rate.",
        "CDC Streams: Tracks real-time telemetry changes on raw tables.",
        "Scheduled Tasks: Runs stored procedures automatically every 60 minutes."
    ], border_color=ACCENT_BLUE)

    # -------------------------------------------------------------
    # SLIDE 6: UiPath Process Automation (RPA)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "Robotic Process Automation (UiPath RPA)")

    add_card(slide6, 0.8, 1.8, 5.6, 5.0, "Automated Workflow Engine", [
        "Directory Watcher: Triggers on new CSV arrival in landing folder.",
        "Database Connector: Invokes Snowflake Stored Procedures securely.",
        "Power BI API Integration: Triggers dataset refresh on completion.",
        "Exception Handling: Try-Catch blocks log operational errors."
    ], border_color=ACCENT_GOLD)

    add_card(slide6, 6.8, 1.8, 5.6, 5.0, "Executive Reporting & Alerts", [
        "Excel Generator: Builds daily fleet executive summary report.",
        "Email Dispatcher: Sends HTML reports to executive stakeholders.",
        "Critical Alert Dispatch: Sends SMS/Email alerts when SOH drops below 80%.",
        "Audit Logging: Writes RPA execution metrics to GOVERNANCE table."
    ], border_color=ACCENT_GOLD)

    # -------------------------------------------------------------
    # SLIDE 7: Power BI & MS Fabric Analytics
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "Power BI & MS Fabric Dark-Blue Dashboards")

    add_card(slide7, 0.8, 1.8, 5.6, 5.0, "6 Interactive Dashboard Pages", [
        "1. Executive Dashboard: Fleet KPIs, cost metrics, active vehicles.",
        "2. Battery Health Analytics: SOH degradation, cell risk, temperature.",
        "3. Charging Operations: Station utilization, peak vs off-peak rates.",
        "4. Fleet Performance: Distance covered, energy efficiency (kWh/100km).",
        "5. Cost Analysis: Energy expense growth, depot tariff breakdown.",
        "6. AI Predictive Dashboard: Maintenance forecasts & AI recommendations."
    ], border_color=ACCENT_CYAN)

    add_card(slide7, 6.8, 1.8, 5.6, 5.0, "Fabric Design & Production DAX", [
        "Styling: Obsidian (#0B132B) & Slate Navy (#1C2541) glassmorphism UI.",
        "DAX SOH Measure: CALCULATE(AVERAGE(FACT[Battery_SOH])).",
        "DAX Cost Growth: DIVIDE([Current_Cost] - [Prior_Cost], [Prior_Cost]).",
        "DirectLake Storage: Real-time query performance over OneLake Delta files."
    ], border_color=ACCENT_CYAN)

    # -------------------------------------------------------------
    # SLIDE 8: AI Assistant & Predictive Engine
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "Interactive AI Insights & Predictive Engine")

    add_card(slide8, 0.8, 1.8, 5.6, 5.0, "Predictive AI Models", [
        "SOH Degradation Predictor: Estimates battery end-of-life timeframe based on charge rate & thermal stress.",
        "Charging Cost Optimizer: Analyzes hourly tariff cycles to recommend shift to off-peak hours (11 PM - 6 AM).",
        "Anomaly Detection: Flags abnormal energy draw or rapid voltage drops."
    ], border_color=ACCENT_GREEN)

    # Embed High-Resolution Battery Health Dashboard Image
    img_path = r"c:\EV DATA DASH\app\battery_dashboard.png"
    if os.path.exists(img_path):
        slide8.shapes.add_picture(img_path, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0))
    else:
        add_card(slide8, 6.8, 1.8, 5.6, 5.0, "Natural Language Console", [
            "Interactive AI Chat Console built directly into the Web Application.",
            "Instant Q&A: Fleet managers query fleet status in plain English.",
            "Example Prompts: 'Which vehicles have SOH < 85%?', 'How to reduce depot 3 charging cost?'",
            "Automated Action Items: Generates maintenance tickets directly."
        ], border_color=ACCENT_GREEN)

    # -------------------------------------------------------------
    # SLIDE 9: Data Governance & Security
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "Enterprise Governance, Security & Lineage")

    add_card(slide9, 0.8, 1.8, 5.6, 5.0, "Security & Compliance", [
        "Role-Based Access (RBAC): EV_ADMIN, EV_DATA_ENGINEER, EV_ANALYST, EV_MAINTENANCE_TECH.",
        "Column-Level Masking: Redacts driver PII & location coordinates for analyst roles.",
        "Encryption: AES-256 bit at rest, TLS 1.3 in transit.",
        "Disaster Recovery: 30-day Snowflake Time Travel & multi-region backup."
    ], border_color=RGBColor(239, 68, 68))

    add_card(slide9, 6.8, 1.8, 5.6, 5.0, "Data Quality & Lineage", [
        "Data Quality Exceptions: Automated rules trap invalid SOH (<50% or >100%) or NULL vehicle IDs.",
        "Audit Logging: All Snowflake transactions & RPA steps logged to GOVERNANCE.AUDIT_LOGS.",
        "Lineage Mapping: Documented flow from Kaggle CSV -> Alteryx -> Snowflake -> Fabric -> AI UI."
    ], border_color=RGBColor(239, 68, 68))

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion & Business Impact
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_header(slide10, "Summary & Next Steps")

    add_card(slide10, 0.8, 1.8, 5.6, 5.0, "Key Project Milestones Achieved", [
        "End-to-End Automated Pipeline operational across Alteryx, Snowflake, UiPath, and Fabric.",
        "Web Application live at http://localhost:3000 with Chart.js & AI Engine.",
        "14% Energy Cost Optimization validated through tariff peak-shifting.",
        "Complete enterprise documentation, ERD diagrams & governance matrices."
    ], border_color=ACCENT_CYAN)

    add_card(slide10, 6.8, 1.8, 5.6, 5.0, "Future Roadmap", [
        "Direct IoT Telemetry Integration via Kafka / AWS IoT Core.",
        "Deep Learning Battery Degradation Forecasting (LSTM models).",
        "Vehicle-to-Grid (V2G) Energy Trading Module.",
        "Mobile App interface for fleet drivers & depot managers."
    ], border_color=ACCENT_CYAN)

    output_path = r"c:\EV DATA DASH\EV_Fleet_Analytics_Platform_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_presentation()
